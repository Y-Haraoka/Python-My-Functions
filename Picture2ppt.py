from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from glob import glob
import os

def Picture2ppt(file, PIC_PER_PAGE, PIC_Size ,FileExtension = 'tif' ,Name=''): 
#file: filepath
#PIC_PER_PAGE:1スライドあたりの画像の数
#PIC_Size: 写真のサイズ（インチ）
#FileExtension: 拡張子

    file_glob = file +"\\*."+FileExtension
    file_names = glob(file_glob)
    File_No = sum(os.path.isfile(os.path.join(file, name)) for name in os.listdir(file))
    IMG_DISPLAY_HEIGHT = Inches(float(PIC_Size)) #画像サイズ
    SLIDES = File_No/PIC_PER_PAGE #スライド数。画像の数と、１スライドあたりの画像の数から算出する。
    
    OUTPUT_FILE_NAME =  Name +".pptx" #出力ファイル名
    
    # Presentaitonインスタンスの作成
    prs = Presentation()
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height

    for i in range(0, File_No, PIC_PER_PAGE): #0からFile_Noまで、PIC_PER_PAGEごとに新しいスライドを作製する。
        blank_slide_layout = prs.slide_layouts[6] #作製するスライドのレイアウト
        slide = prs.slides.add_slide(blank_slide_layout) #スライドの追加
        
        for j in range(0, PIC_PER_PAGE): #0～PIC_PER_PAGEまでの画像を1つのスライドに貼り付ける

						#貼り付ける画像ファイル名を取得
            file_name = file_names[i+j] 
            Figure_name = os.path.basename(file_name) #画像の上に貼り付けるファイル名として使用する。
            
            #画像サイズを取得してアスペクト比を得る。そのままの画像だと大き過ぎるので。
            im = Image.open(file_name) #選択した画像を開く
            im_width, im_height = im.size #.sizeはタプル形式で（幅、高さ）を返すので、変数を2つ(im_width, im_height)を準備して、そこに数値を格納する
            aspect_ratio = im_width/im_height
            
            #表示された画像のサイズを計算
            img_display_height = IMG_DISPLAY_HEIGHT #画像の高さは最初に設定したインチ
            img_display_width = aspect_ratio*img_display_height #アスペクト比をもとにして横の長さを規定する
            
#画像の位置を設定する
            left = (( SLIDE_WIDTH / PIC_PER_PAGE - img_display_width ) / PIC_PER_PAGE ) + SLIDE_WIDTH *j/ PIC_PER_PAGE
            top = ( SLIDE_HEIGHT - img_display_height ) /PIC_PER_PAGE

            slide.shapes.add_picture(file_name, left, top, height=IMG_DISPLAY_HEIGHT)
            
            #画像名テキストを追加
            width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top*0.5, width*2.5, height)
            tf = txBox.text_frame
            tf.word_wrap = True
                                
            tf.text = Figure_name.replace(FileExtension,'')
                            

    return prs.save(file +OUTPUT_FILE_NAME)